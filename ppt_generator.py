"""Generate PowerPoint presentations from metrics"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from typing import Dict
from datetime import datetime
import matplotlib.pyplot as plt
import matplotlib
matplotlib.use('Agg')  # Use non-interactive backend
import io


class PPTGenerator:
    """Generate PowerPoint presentations with metrics"""
    
    def __init__(self):
        """Initialize presentation"""
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
    
    def create_title_slide(self, team_name: str, sprint_name: str):
        """Create title slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = f"{team_name} Sprint Velocity Report"
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(44)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(0, 51, 102)
        title_paragraph.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = f"Sprint: {sprint_name} | Generated: {datetime.now().strftime('%B %d, %Y')}"
        subtitle_paragraph = subtitle_frame.paragraphs[0]
        subtitle_paragraph.font.size = Pt(18)
        subtitle_paragraph.font.color.rgb = RGBColor(64, 64, 64)
        subtitle_paragraph.alignment = PP_ALIGN.CENTER
    
    def create_current_sprint_slide(self, metrics: Dict):
        """Create slide with current sprint metrics"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "Current Sprint Metrics"
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(32)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(0, 51, 102)
        
        current = metrics.get('current_sprint', {})
        y_pos = 1.5
        
        # Story Points Committed
        self._add_metric_box(
            slide, "Story Points Committed", 
            current.get('committed_story_points', 0),
            Inches(0.5), Inches(y_pos), Inches(4), Inches(1.2)
        )
        
        # Story Points Completed
        self._add_metric_box(
            slide, "Story Points Completed",
            current.get('completed_story_points', 0),
            Inches(5.5), Inches(y_pos), Inches(4), Inches(1.2)
        )
        
        y_pos += 1.5
        
        # Completion Rate
        completion_rate = current.get('completion_rate', 0)
        self._add_metric_box(
            slide, "Completion Rate",
            f"{completion_rate}%",
            Inches(0.5), Inches(y_pos), Inches(4), Inches(1.2)
        )
        
        # Defect Count
        self._add_metric_box(
            slide, "Defects Found",
            current.get('defect_count', 0),
            Inches(5.5), Inches(y_pos), Inches(4), Inches(1.2)
        )
        
        y_pos += 1.5
        
        # Total Issues
        self._add_metric_box(
            slide, "Total Issues",
            current.get('total_issues', 0),
            Inches(0.5), Inches(y_pos), Inches(4), Inches(1.2)
        )
        
        # Completed Issues
        self._add_metric_box(
            slide, "Completed Issues",
            current.get('completed_issues', 0),
            Inches(5.5), Inches(y_pos), Inches(4), Inches(1.2)
        )
    
    def create_velocity_improvement_slide(self, metrics: Dict):
        """Create slide showing velocity improvement"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "Velocity Improvement After AI Adoption"
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(32)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(0, 51, 102)
        
        improvement = metrics.get('velocity_improvement', {})
        baseline = metrics.get('baseline_velocity', {})
        post_ai = metrics.get('post_ai_velocity', {})
        
        y_pos = 1.5
        
        # Baseline Velocity
        self._add_metric_box(
            slide, f"Baseline Velocity\n(Before AI)",
            f"{improvement.get('baseline_velocity', 0)} SP",
            Inches(0.5), Inches(y_pos), Inches(3), Inches(1.5)
        )
        
        # Post AI Velocity
        self._add_metric_box(
            slide, f"Post-AI Velocity",
            f"{improvement.get('post_ai_velocity', 0)} SP",
            Inches(3.75), Inches(y_pos), Inches(3), Inches(1.5)
        )
        
        # Improvement
        improvement_percent = improvement.get('improvement_percent', 0)
        improvement_color = RGBColor(0, 128, 0) if improvement_percent > 0 else RGBColor(128, 0, 0)
        
        self._add_metric_box(
            slide, f"Improvement",
            f"{improvement_percent}%",
            Inches(7), Inches(y_pos), Inches(2.5), Inches(1.5),
            value_color=improvement_color
        )
        
        # Chart
        y_pos += 2
        self._add_velocity_chart(slide, baseline, post_ai, improvement, Inches(1), Inches(y_pos), Inches(8), Inches(3))
    
    def create_defect_metrics_slide(self, metrics: Dict):
        """Create slide with defect metrics"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "Defect Metrics"
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(32)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(0, 51, 102)
        
        defect_metrics = metrics.get('defect_metrics', {})
        y_pos = 1.5
        
        # Baseline Defects
        self._add_metric_box(
            slide, f"Avg Defects\n(Before AI)",
            f"{defect_metrics.get('baseline_avg_defects', 0)}",
            Inches(0.5), Inches(y_pos), Inches(3), Inches(1.5)
        )
        
        # Post AI Defects
        self._add_metric_box(
            slide, f"Avg Defects\n(After AI)",
            f"{defect_metrics.get('post_ai_avg_defects', 0)}",
            Inches(3.75), Inches(y_pos), Inches(3), Inches(1.5)
        )
        
        # Reduction
        reduction_percent = defect_metrics.get('defect_reduction_percent', 0)
        reduction_color = RGBColor(0, 128, 0) if reduction_percent > 0 else RGBColor(128, 0, 0)
        
        self._add_metric_box(
            slide, f"Defect Reduction",
            f"{reduction_percent}%",
            Inches(7), Inches(y_pos), Inches(2.5), Inches(1.5),
            value_color=reduction_color
        )
        
        # Chart
        y_pos += 2
        self._add_defect_chart(slide, defect_metrics, Inches(1), Inches(y_pos), Inches(8), Inches(3))
    
    def create_summary_slide(self, metrics: Dict):
        """Create summary slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "Key Takeaways"
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(32)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(0, 51, 102)
        
        current = metrics.get('current_sprint', {})
        improvement = metrics.get('velocity_improvement', {})
        defect_metrics = metrics.get('defect_metrics', {})
        
        # Summary points
        summary_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
        summary_frame = summary_box.text_frame
        summary_frame.word_wrap = True
        
        points = [
            f"• Story Points Committed: {current.get('committed_story_points', 0)}",
            f"• Velocity Improvement: {improvement.get('improvement_percent', 0)}%",
            f"• Defect Reduction: {defect_metrics.get('defect_reduction_percent', 0)}%",
            f"• AI Adoption Date: {metrics.get('ai_adoption_date', 'N/A')}",
        ]
        
        for i, point in enumerate(points):
            p = summary_frame.add_paragraph()
            p.text = point
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(64, 64, 64)
            p.space_after = Pt(12)
            if i == 0:
                p.level = 0
    
    def _add_metric_box(self, slide, label: str, value, left, top, width, height, value_color=None):
        """Add a metric box to slide"""
        box = slide.shapes.add_textbox(left, top, width, height)
        frame = box.text_frame
        frame.text = f"{label}\n\n{value}"
        frame.word_wrap = True
        
        # Format label
        label_para = frame.paragraphs[0]
        label_para.font.size = Pt(14)
        label_para.font.bold = True
        label_para.font.color.rgb = RGBColor(64, 64, 64)
        label_para.alignment = PP_ALIGN.CENTER
        
        # Format value
        if len(frame.paragraphs) > 1:
            value_para = frame.paragraphs[1]
            value_para.font.size = Pt(28)
            value_para.font.bold = True
            value_para.font.color.rgb = value_color or RGBColor(0, 51, 102)
            value_para.alignment = PP_ALIGN.CENTER
        
        # Add border
        box.line.color.rgb = RGBColor(200, 200, 200)
        box.line.width = Pt(1)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(245, 245, 245)
    
    def _add_velocity_chart(self, slide, baseline, post_ai, improvement, left, top, width, height):
        """Add velocity comparison chart"""
        fig, ax = plt.subplots(figsize=(width.inches, height.inches))
        
        categories = ['Baseline\n(Before AI)', 'Post-AI']
        velocities = [
            improvement.get('baseline_velocity', 0),
            improvement.get('post_ai_velocity', 0)
        ]
        
        colors = ['#FF6B6B', '#4ECDC4']
        bars = ax.bar(categories, velocities, color=colors, alpha=0.8)
        
        ax.set_ylabel('Story Points', fontsize=12, fontweight='bold')
        ax.set_title('Velocity Comparison', fontsize=14, fontweight='bold', pad=20)
        ax.grid(axis='y', alpha=0.3)
        
        # Add value labels on bars
        for bar in bars:
            height = bar.get_height()
            ax.text(bar.get_x() + bar.get_width()/2., height,
                   f'{height:.1f}',
                   ha='center', va='bottom', fontsize=12, fontweight='bold')
        
        plt.tight_layout()
        
        # Save to buffer
        buf = io.BytesIO()
        plt.savefig(buf, format='png', dpi=150, bbox_inches='tight')
        buf.seek(0)
        plt.close()
        
        # Add to slide
        slide.shapes.add_picture(buf, left, top, width, height)
    
    def _add_defect_chart(self, slide, defect_metrics, left, top, width, height):
        """Add defect comparison chart"""
        fig, ax = plt.subplots(figsize=(width.inches, height.inches))
        
        categories = ['Baseline\n(Before AI)', 'Post-AI']
        defects = [
            defect_metrics.get('baseline_avg_defects', 0),
            defect_metrics.get('post_ai_avg_defects', 0)
        ]
        
        colors = ['#FF6B6B', '#4ECDC4']
        bars = ax.bar(categories, defects, color=colors, alpha=0.8)
        
        ax.set_ylabel('Average Defects', fontsize=12, fontweight='bold')
        ax.set_title('Defect Comparison', fontsize=14, fontweight='bold', pad=20)
        ax.grid(axis='y', alpha=0.3)
        
        # Add value labels on bars
        for bar in bars:
            height = bar.get_height()
            ax.text(bar.get_x() + bar.get_width()/2., height,
                   f'{height:.1f}',
                   ha='center', va='bottom', fontsize=12, fontweight='bold')
        
        plt.tight_layout()
        
        # Save to buffer
        buf = io.BytesIO()
        plt.savefig(buf, format='png', dpi=150, bbox_inches='tight')
        buf.seek(0)
        plt.close()
        
        # Add to slide
        slide.shapes.add_picture(buf, left, top, width, height)
    
    def save(self, filename: str):
        """Save presentation to file"""
        self.prs.save(filename)
        print(f"Presentation saved to {filename}")
    
    def create_ai_impact_slide(self, metrics: Dict):
        """Create slide showing AI story points comparison"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "AI Impact: Time Saved Analysis"
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(32)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(0, 51, 102)
        
        current = metrics.get('current_sprint', {})
        
        # Only show if AI data is available
        if not current.get('has_ai_data', False):
            no_data_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
            no_data_frame = no_data_box.text_frame
            no_data_frame.text = "AI Story Points data not available.\nPlease configure AI_STORY_POINTS_FIELD_ID in .env file."
            no_data_paragraph = no_data_frame.paragraphs[0]
            no_data_paragraph.font.size = Pt(18)
            no_data_paragraph.font.color.rgb = RGBColor(128, 128, 128)
            no_data_paragraph.alignment = PP_ALIGN.CENTER
            return
        
        y_pos = 1.5
        
        # AI Story Points (without AI)
        self._add_metric_box(
            slide, "Estimated Story Points\n(Without AI)",
            f"{current.get('ai_story_points_committed', 0)} SP",
            Inches(0.5), Inches(y_pos), Inches(3), Inches(1.5)
        )
        
        # Actual Story Points (with AI)
        self._add_metric_box(
            slide, "Actual Story Points\n(With AI)",
            f"{current.get('committed_story_points', 0)} SP",
            Inches(3.75), Inches(y_pos), Inches(3), Inches(1.5)
        )
        
        # Time Saved
        time_saved = current.get('time_saved_total', 0)
        time_saved_percent = current.get('time_saved_percent', 0)
        self._add_metric_box(
            slide, "Time Saved",
            f"{time_saved} SP\n({time_saved_percent}%)",
            Inches(7), Inches(y_pos), Inches(2.5), Inches(1.5),
            RGBColor(0, 128, 0) if time_saved > 0 else RGBColor(128, 128, 128)
        )
        
        y_pos += 2
        
        # Chart comparing AI vs Actual
        self._add_ai_comparison_chart(
            slide, current,
            Inches(1), Inches(y_pos), Inches(8), Inches(3.5)
        )
    
    def _add_ai_comparison_chart(self, slide, current_metrics, left, top, width, height):
        """Add AI vs Actual story points comparison chart"""
        fig, ax = plt.subplots(figsize=(width.inches, height.inches))
        
        categories = ['Estimated\n(Without AI)', 'Actual\n(With AI)']
        values = [
            current_metrics.get('ai_story_points_committed', 0),
            current_metrics.get('committed_story_points', 0)
        ]
        
        colors = ['#FF6B6B', '#4ECDC4']
        bars = ax.bar(categories, values, color=colors, alpha=0.8)
        
        ax.set_ylabel('Story Points', fontsize=12, fontweight='bold')
        ax.set_title('AI Impact: Story Points Comparison', fontsize=14, fontweight='bold', pad=20)
        ax.grid(axis='y', alpha=0.3)
        
        # Add value labels on bars
        for bar in bars:
            bar_height = bar.get_height()
            ax.text(bar.get_x() + bar.get_width()/2., bar_height,
                   f'{bar_height:.1f} SP',
                   ha='center', va='bottom', fontsize=12, fontweight='bold')
        
        # Add time saved annotation
        time_saved = current_metrics.get('time_saved_total', 0)
        if time_saved > 0:
            ax.text(0.5, max(values) * 0.9, 
                   f'Time Saved: {time_saved} SP ({current_metrics.get("time_saved_percent", 0)}%)',
                   ha='center', va='center', fontsize=11, fontweight='bold',
                   bbox=dict(boxstyle='round', facecolor='lightgreen', alpha=0.7))
        
        plt.tight_layout()
        
        # Save to buffer
        buf = io.BytesIO()
        plt.savefig(buf, format='png', dpi=150, bbox_inches='tight')
        buf.seek(0)
        plt.close()
        
        # Add to slide
        slide.shapes.add_picture(buf, left, top, width, height)
    
    def generate_presentation(self, team_name: str, metrics: Dict, output_file: str):
        """Generate complete presentation"""
        sprint_name = metrics.get('current_sprint', {}).get('sprint_name', 'Current Sprint')
        
        self.create_title_slide(team_name, sprint_name)
        self.create_current_sprint_slide(metrics)
        
        # Add AI impact slide if data is available
        current = metrics.get('current_sprint', {})
        if current.get('has_ai_data', False):
            self.create_ai_impact_slide(metrics)
        
        self.create_velocity_improvement_slide(metrics)
        self.create_defect_metrics_slide(metrics)
        self.create_summary_slide(metrics)
        
        self.save(output_file)
